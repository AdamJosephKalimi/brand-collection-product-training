"""
Presentation Generation Service

Generates PowerPoint presentations from collection data including:
- Intro slides (LLM-generated content)
- Product slides (collection items with dynamic layout)
"""

import logging
import os
import tempfile
from collections import OrderedDict
from datetime import datetime
from typing import Dict, Any, Optional
from io import BytesIO

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
from pptx.dml.color import RGBColor

from app.services.firebase_service import firebase_service
from app.services.collection_service import collection_service
from app.services.llm_service import llm_service
from app.services.intro_slide_generation_service import LANGUAGE_NAMES

logger = logging.getLogger(__name__)


class PresentationGenerationService:
    """Service for generating PowerPoint presentations from collection data"""
    
    def __init__(self):
        self.db = firebase_service.db
        self.bucket = firebase_service.bucket
        self.prs = None
        self.blank_layout = None
        self.title_slide_layout = None
        self.title_content_layout = None
        self._w_scale = 1.0  # Horizontal scale factor (1.0 for 4:3, 1.333 for 16:9)
        self._h_offset = 0  # Horizontal offset in inches to center 10"-wide content in wider slides
        self._typo = None  # Deck typography settings dict
        self._language = "en"  # Deck output language
        self._translations = {}  # {original_string: translated_string}
    
    async def generate_presentation(
        self,
        collection_id: str,
        user_id: str,
        products_per_slide: int = 1,
        slide_aspect_ratio: str = "16:9",
        deck_typography: dict = None,
        selected_language: str = "en"
    ) -> str:
        """
        Generate complete PowerPoint presentation for a collection.

        Args:
            collection_id: ID of the collection
            user_id: ID of the user requesting generation
            products_per_slide: Number of products per slide (1, 2, or 4)
            slide_aspect_ratio: Slide aspect ratio ("4:3" or "16:9")
            deck_typography: Optional dict with heading/body text style overrides
            selected_language: Language code for translatable deck content (default "en")

        Returns:
            Download URL for the generated presentation
        """
        try:
            logger.info(f"Starting presentation generation for collection: {collection_id}")

            self._language = selected_language
            self._translations = {}

            # 1. Fetch collection data
            collection = await collection_service.get_collection(collection_id, user_id)
            
            if not collection:
                raise ValueError(f"Collection not found: {collection_id}")
            
            intro_slides = collection.intro_slides

            # Extract field visibility settings for product slides
            s = collection.settings
            self._show_fields = {
                'product_name': getattr(s, 'show_product_name', True),
                'sku': getattr(s, 'show_sku', True),
                'description': getattr(s, 'show_descriptions', True),
                'color': getattr(s, 'show_color', True),
                'material': getattr(s, 'show_material', True),
                'sizes': getattr(s, 'show_sizes', True),
                'origin': getattr(s, 'show_origin', True),
                'wholesale_price': getattr(s, 'show_wholesale_price', False),
                'rrp': getattr(s, 'show_rrp', True),
            }

            if not intro_slides or not intro_slides.get('slides'):
                raise ValueError(f"No intro slides found for collection: {collection_id}")
            
            logger.info(f"Found {len(intro_slides['slides'])} intro slides")
            
            # 2. Initialize presentation
            self.prs = Presentation()

            # Set slide dimensions based on aspect ratio
            if slide_aspect_ratio == "16:9":
                self.prs.slide_width = Inches(13.333)
                self.prs.slide_height = Inches(7.5)
                self._w_scale = 13.333 / 10.0  # ~1.333
                self._h_offset = (13.333 - 10.0) / 2.0  # ~1.667" per side
            else:
                # 4:3 (default python-pptx dimensions)
                self.prs.slide_width = Inches(10)
                self.prs.slide_height = Inches(7.5)
                self._w_scale = 1.0
                self._h_offset = 0

            self._typo = deck_typography or {}
            if self._typo:
                logger.info(f"Deck typography loaded: {self._typo}")

            self.blank_layout = self.prs.slide_layouts[6]  # Blank layout
            self.title_slide_layout = self.prs.slide_layouts[0]  # Title Slide
            self.title_content_layout = self.prs.slide_layouts[1]  # Title and Content

            # Center intro layout placeholders for widescreen (no-op for 4:3)
            self._center_layout_placeholders()

            # 3. Generate intro slides
            await self._generate_intro_slides(intro_slides)
            
            # 4. Generate product slides
            category_order = {}
            for cat in (collection.categories or []):
                cat_dict = cat if isinstance(cat, dict) else cat.dict()
                category_order[cat_dict['name']] = cat_dict.get('display_order', 0)
            await self._generate_product_slides(collection_id, products_per_slide, category_order)
            
            # 5. Save and upload
            download_url = await self._save_and_upload(collection_id)
            
            logger.info(f"✅ Presentation generated successfully: {download_url}")
            
            return download_url
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise
    
    async def _generate_intro_slides(self, intro_slides: Dict[str, Any]):
        """
        Generate intro slides from intro_slides data.
        Currently supports: cover_page, brand_history
        
        Args:
            intro_slides: Dictionary containing slides array
        """
        logger.info("Generating intro slides...")
        
        slides_generated = 0
        
        for slide_data in intro_slides['slides']:
            slide_type = slide_data.get('slide_type')
            
            if slide_type == 'cover_page':
                logger.info("Creating cover page slide...")
                self._create_cover_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'brand_history':
                logger.info("Creating brand history slide...")
                self._create_brand_history_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'brand_introduction':
                logger.info("Creating brand introduction slide...")
                self._create_brand_introduction_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'brand_values':
                logger.info("Creating brand values slide...")
                self._create_brand_values_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'brand_personality':
                logger.info("Creating brand personality slide...")
                self._create_brand_personality_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'flagship_stores':
                logger.info("Creating flagship stores slide...")
                self._create_flagship_stores_slide(slide_data)
                slides_generated += 1
                
            elif slide_type == 'collection_introduction':
                logger.info("Creating collection introduction slide...")
                self._create_collection_introduction_slide(slide_data)
                slides_generated += 1

            elif slide_type == 'core_collection':
                logger.info("Creating core collection slide...")
                self._create_core_collection_slide(slide_data)
                slides_generated += 1
            
            else:
                logger.debug(f"Unknown slide type: {slide_type}")
        
        logger.info(f"Generated {slides_generated} slides")
    
    def _center_layout_placeholders(self):
        """
        Center title and content placeholders in the intro slide layouts
        for widescreen. Modifies the LAYOUT-level placeholders so that
        slide-level shapes inherit correct centered positions.

        CRITICAL: Must read ALL four geometry properties (left, top, width,
        height) BEFORE writing ANY of them. Setting just `ph.left` on a
        placeholder that inherits its xfrm from the slide master creates a
        new xfrm element with zeroed width/height, making it invisible.
        """
        if self._h_offset == 0:
            return
        slide_w = self.prs.slide_width
        for layout in [self.title_slide_layout, self.title_content_layout]:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx in (0, 1):
                    # Read ALL inherited values BEFORE any setter triggers xfrm creation
                    w = ph.width
                    h = ph.height
                    t = ph.top
                    # Write all four to create a complete xfrm element
                    ph.left = (slide_w - w) // 2
                    ph.top = t
                    ph.width = w
                    ph.height = h

    def _parse_hex_color(self, hex_str: str) -> RGBColor:
        """Parse a hex color string like '#2C3528' into an RGBColor."""
        hex_str = hex_str.lstrip('#')
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    def _apply_typo(self, font, group: str, default_size: int = None):
        """
        Apply deck typography overrides to a font object.

        Args:
            font: pptx font object (p.font or run.font)
            group: 'heading' or 'body'
            default_size: default Pt size to use if no override
        """
        style = self._typo.get(group) if self._typo else None
        if style:
            if style.get('font_family'):
                font.name = style['font_family']
            if style.get('font_size'):
                font.size = Pt(style['font_size'])
            elif default_size:
                font.size = Pt(default_size)
            if style.get('font_color'):
                try:
                    font.color.rgb = self._parse_hex_color(style['font_color'])
                except (ValueError, IndexError):
                    pass
        elif default_size:
            font.size = Pt(default_size)

    def _apply_body_font(self, text_frame):
        """
        Apply body typography font family and color to all paragraphs in a text frame.
        Preserves existing font sizes. Call after all text has been added.
        """
        style = self._typo.get('body') if self._typo else None
        if not style:
            return
        for p in text_frame.paragraphs:
            for run in p.runs:
                if style.get('font_family'):
                    run.font.name = style['font_family']
                if style.get('font_color'):
                    try:
                        run.font.color.rgb = self._parse_hex_color(style['font_color'])
                    except (ValueError, IndexError):
                        pass
            # Also apply to paragraph-level font (for text set via p.text)
            if style.get('font_family'):
                p.font.name = style['font_family']
            if style.get('font_color'):
                try:
                    p.font.color.rgb = self._parse_hex_color(style['font_color'])
                except (ValueError, IndexError):
                    pass

    def _t(self, text: str) -> str:
        """Look up a translated string, falling back to the original."""
        if not text or self._language == "en":
            return text
        return self._translations.get(text, text)

    async def _translate_strings(self, strings: list) -> dict:
        """
        Batch-translate a list of strings via LLM.

        Args:
            strings: List of unique strings to translate

        Returns:
            Dict mapping original string -> translated string
        """
        if not strings or self._language == "en":
            return {}

        lang_name = LANGUAGE_NAMES.get(self._language, self._language)
        # Deduplicate while preserving order
        unique = list(dict.fromkeys(s for s in strings if s))
        if not unique:
            return {}

        import json as _json
        prompt = f"""Translate the following strings from English to {lang_name}.
Return ONLY a JSON object mapping each original English string to its {lang_name} translation.
Do NOT translate proper nouns (brand names, product names, city names).
Keep the translations concise and natural.

Strings to translate:
{_json.dumps(unique, ensure_ascii=False)}

Return ONLY valid JSON like:
{{
    "Original string": "Translated string"
}}"""

        try:
            result = await llm_service.generate_json_completion(
                prompt=prompt,
                temperature=0.3,
                max_tokens=1000
            )
            if isinstance(result, dict):
                logger.info(f"Translated {len(result)} strings to {lang_name}")
                return result
        except Exception as e:
            logger.warning(f"Batch translation failed: {e}")
        return {}

    def _create_cover_slide(self, data: Dict[str, Any]):
        """
        Create cover page slide using Title Slide layout.
        
        Uses PowerPoint's built-in Title Slide layout with:
        - Title placeholder for collection title
        - Subtitle placeholder for brand name
        - Manual textbox for tagline (if exists)
        
        Args:
            data: Slide data containing title, subtitle, content
        """
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        content = data.get('content', {})

        # Use title placeholder
        title = slide.shapes.title
        title.text = content.get('title', data.get('title', 'Collection'))
        
        # Format title
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.font.bold = True
        self._apply_typo(p.font, 'heading', default_size=44)

        # Use subtitle placeholder (usually placeholders[1] in title slide)
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = content.get('subtitle', data.get('subtitle', ''))

            # Format subtitle
            tf = subtitle.text_frame
            p = tf.paragraphs[0]
            self._apply_typo(p.font, 'heading', default_size=24)
        
        # Tagline - add as manual textbox at bottom, centered across full slide width
        if content.get('tagline'):
            tagline_width = 8.0  # inches
            slide_w = 13.333 if self._w_scale > 1 else 10.0
            tagline_left = (slide_w - tagline_width) / 2.0
            tagline_box = slide.shapes.add_textbox(
                left=Inches(tagline_left),
                top=Inches(6.5),
                width=Inches(tagline_width),
                height=Inches(0.5)
            )
            tf = tagline_box.text_frame
            tf.text = content.get('tagline')
            
            # Format tagline
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            self._apply_body_font(tf)

        logger.info("Cover slide created (using Title Slide layout)")
    
    def _create_brand_history_slide(self, data: Dict[str, Any]):
        """
        Create brand history slide using Title and Content layout.
        
        Args:
            data: Slide data containing title, content with founding, evolution, market_context, philosophical_takeaway
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Brand History')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Extract nested founding structure
        founding = content.get('founding', {})
        if founding and isinstance(founding, dict):
            year = founding.get('year', '')
            founder_info = founding.get('founder', {})
            origin_info = founding.get('origin', {})
            
            # Founder name and background
            founder_name = founder_info.get('name', '') if isinstance(founder_info, dict) else str(founder_info) if founder_info else ''
            founder_bg = founder_info.get('background', '') if isinstance(founder_info, dict) else ''
            
            # Origin city and country
            city = origin_info.get('city', '') if isinstance(origin_info, dict) else str(origin_info) if origin_info else ''
            country = origin_info.get('country', '') if isinstance(origin_info, dict) else ''
            
            # Build founding text
            if year or founder_name:
                p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                first_paragraph = False
                
                founding_text = f"Founded in {year}" if year else "Founded"
                if founder_name:
                    founding_text += f" by {founder_name}"
                if city or country:
                    location = f"{city}, {country}" if city and country else city or country
                    founding_text += f" in {location}"
                
                p.text = founding_text
                p.font.size = Pt(14)
                p.font.bold = True
                p.level = 0
            
            if founder_bg:
                p = tf.add_paragraph()
                p.text = founder_bg
                p.font.size = Pt(13)
                p.level = 0
        
        # Fallback for flat structure (backward compatibility)
        elif content.get('founding_year') or content.get('founder'):
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            p.text = f"Founded: {content.get('founding_year', '')}"
            if content.get('founder'):
                p.text += f" by {content.get('founder')}"
            p.font.size = Pt(14)
            p.level = 0
        
        # Add evolution section
        evolution = content.get('evolution', {})
        if evolution and isinstance(evolution, dict):
            early_identity = evolution.get('early_identity', '')
            key_shifts = evolution.get('key_shifts', [])
            signature_elements = evolution.get('signature_elements', [])
            
            if early_identity or key_shifts or signature_elements:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                p = tf.add_paragraph()
                p.text = "Brand Evolution:"
                p.font.size = Pt(14)
                p.font.bold = True
                p.level = 0
                
                if early_identity:
                    p = tf.add_paragraph()
                    p.text = early_identity
                    p.font.size = Pt(12)
                    p.level = 1
                
                for shift in key_shifts:
                    p = tf.add_paragraph()
                    p.text = shift
                    p.font.size = Pt(12)
                    p.level = 1
                
                if signature_elements:
                    p = tf.add_paragraph()
                    p.text = ""  # Spacing
                    
                    p = tf.add_paragraph()
                    p.text = "Signature Elements:"
                    p.font.size = Pt(13)
                    p.font.bold = True
                    p.level = 0
                    
                    for element in signature_elements:
                        p = tf.add_paragraph()
                        p.text = element
                        p.font.size = Pt(12)
                        p.level = 1
        
        # Add market context
        market_context = content.get('market_context', {})
        if market_context and isinstance(market_context, dict):
            positioning = market_context.get('positioning', '')
            production = market_context.get('production_philosophy', '')
            footprint = market_context.get('current_footprint', '')
            
            if positioning or production or footprint:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                if positioning:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = "Market Position: "
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run = p.add_run()
                    run.text = positioning
                    run.font.size = Pt(12)
                    p.level = 0
                
                if production:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = "Production: "
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run = p.add_run()
                    run.text = production
                    run.font.size = Pt(12)
                    p.level = 0
        
        # Add philosophical takeaway
        takeaway = content.get('philosophical_takeaway', '')
        if takeaway:
            p = tf.add_paragraph()
            p.text = ""  # Spacing
            
            p = tf.add_paragraph()
            p.text = takeaway
            p.font.size = Pt(12)
            p.font.italic = True
            p.level = 0
        
        self._apply_body_font(tf)
        logger.info("Brand history slide created (using Title and Content layout)")
    
    def _create_brand_introduction_slide(self, data: Dict[str, Any]):
        """
        Create brand introduction slide using Title and Content layout.
        
        Args:
            data: Slide data containing title, content with overview, brand_identity, name_and_history
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Brand Introduction')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Add overview section (nested object with summary and positioning)
        overview = content.get('overview', {})
        if isinstance(overview, dict):
            summary = overview.get('summary', '')
            positioning = overview.get('positioning', '')
            
            if summary:
                p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                first_paragraph = False
                p.text = summary
                p.font.size = Pt(14)
                p.level = 0
            
            if positioning:
                p = tf.add_paragraph()
                p.text = positioning
                p.font.size = Pt(14)
                p.level = 0
        elif overview:  # Fallback if it's a string
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            p.text = str(overview)
            p.font.size = Pt(14)
            p.level = 0
        
        # Add brand identity section
        brand_identity = content.get('brand_identity', {})
        if brand_identity:
            essence = brand_identity.get('essence', '')
            differentiators = brand_identity.get('differentiators', [])
            unique_facts = brand_identity.get('unique_facts', [])
            
            if essence:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Brand Essence: "
                run.font.bold = True
                run.font.size = Pt(14)
                run = p.add_run()
                run.text = essence
                run.font.size = Pt(14)
                p.level = 0
            
            if differentiators:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                p = tf.add_paragraph()
                p.text = "What Sets Us Apart:"
                p.font.bold = True
                p.font.size = Pt(14)
                p.level = 0
                
                for diff in differentiators:
                    p = tf.add_paragraph()
                    p.text = diff
                    p.font.size = Pt(13)
                    p.level = 1
            
            if unique_facts:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                p = tf.add_paragraph()
                p.text = "Did You Know:"
                p.font.bold = True
                p.font.size = Pt(14)
                p.level = 0
                
                for fact in unique_facts:
                    p = tf.add_paragraph()
                    p.text = fact
                    p.font.size = Pt(13)
                    p.level = 1
        
        # Add name and history section if present
        name_history = content.get('name_and_history', {})
        if name_history:
            name_origin = name_history.get('name_origin', '')
            background = name_history.get('background', '')
            
            if name_origin or background:
                p = tf.add_paragraph()
                p.text = ""  # Spacing
                
                if name_origin:
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = "Name Origin: "
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run = p.add_run()
                    run.text = name_origin
                    run.font.size = Pt(14)
                    p.level = 0
                
                if background:
                    p = tf.add_paragraph()
                    p.text = background
                    p.font.size = Pt(13)
                    p.level = 0
        
        self._apply_body_font(tf)
        logger.info("Brand introduction slide created (using Title and Content layout)")
    
    def _create_brand_values_slide(self, data: Dict[str, Any]):
        """
        Create brand values slide using Title and Content layout.
        
        Args:
            data: Slide data containing title, content with values array (headline + explanation)
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Brand Values')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        # Add values - LLM returns {headline, explanation} for each value
        values = content.get('values', [])
        if values:
            p = tf.paragraphs[0]
            first = True
            
            for value in values:
                if not first:
                    p = tf.add_paragraph()
                first = False
                
                # Value headline and explanation (also support name/description as fallback)
                value_name = value.get('headline', value.get('name', '')) if isinstance(value, dict) else str(value)
                value_desc = value.get('explanation', value.get('description', '')) if isinstance(value, dict) else ''
                
                p.level = 0
                
                if value_desc:
                    # Add value name in bold
                    run = p.add_run()
                    run.text = f"{value_name}: "
                    run.font.bold = True
                    run.font.size = Pt(14)
                    
                    # Add explanation in normal text
                    run = p.add_run()
                    run.text = value_desc
                    run.font.size = Pt(14)
                else:
                    # Just the value name
                    p.text = value_name
                    p.font.size = Pt(14)
        
        self._apply_body_font(tf)
        logger.info("Brand values slide created (using Title and Content layout)")
    
    def _create_brand_personality_slide(self, data: Dict[str, Any]):
        """
        Create brand personality slide using Title and Content layout.
        
        Args:
            data: Slide data containing title, content with quote and traits (name + explanation)
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Brand Personality')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Add quote section if present
        quote = content.get('quote', {})
        if quote and isinstance(quote, dict):
            quote_text = quote.get('text', '')
            attribution = quote.get('attribution', '')
            
            if quote_text:
                p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                first_paragraph = False
                p.text = f'"{quote_text}"'
                p.font.size = Pt(14)
                p.font.italic = True
                p.level = 0
                
                if attribution:
                    p = tf.add_paragraph()
                    p.text = f"— {attribution}"
                    p.font.size = Pt(12)
                    p.level = 0
                
                # Add spacing after quote
                p = tf.add_paragraph()
                p.text = ""
        
        # Add traits section - LLM returns {name, explanation} for each trait
        traits = content.get('traits', content.get('personality_traits', []))
        if traits:
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            
            p.text = "Personality Traits:"
            p.font.size = Pt(14)
            p.font.bold = True
            p.level = 0
            
            for trait in traits:
                p = tf.add_paragraph()
                
                if isinstance(trait, dict):
                    trait_name = trait.get('name', '')
                    explanation = trait.get('explanation', '')
                    
                    if explanation:
                        run = p.add_run()
                        run.text = f"{trait_name}: "
                        run.font.bold = True
                        run.font.size = Pt(13)
                        
                        run = p.add_run()
                        run.text = explanation
                        run.font.size = Pt(13)
                    else:
                        p.text = trait_name
                        p.font.size = Pt(13)
                else:
                    p.text = str(trait)
                    p.font.size = Pt(13)
                
                p.level = 1
        
        self._apply_body_font(tf)
        logger.info("Brand personality slide created (using Title and Content layout)")
    
    def _create_flagship_stores_slide(self, data: Dict[str, Any]):
        """
        Create flagship stores slide using Title and Content layout.
        
        Args:
            data: Slide data containing title, content with flagship_stores array and alternative_experiences
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Flagship Stores & Experiences')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Add flagship stores - LLM returns nested structure with name, location, design_elements, etc.
        stores = content.get('flagship_stores', content.get('store_locations', content.get('stores', [])))
        alternative = content.get('alternative_experiences', '')
        
        if stores:
            for store in stores:
                if isinstance(store, dict):
                    # Get store name and location
                    store_name = store.get('name', '')
                    location = store.get('location', {})
                    city = location.get('city', '') if isinstance(location, dict) else str(location)
                    country = location.get('country', '') if isinstance(location, dict) else ''
                    year_opened = store.get('year_opened', '')
                    
                    # Build store header
                    location_str = f"{city}, {country}" if city and country else city or country
                    header = store_name or location_str
                    if year_opened and header:
                        header = f"{header} ({year_opened})"
                    
                    if header:
                        p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                        first_paragraph = False
                        p.text = header
                        p.font.size = Pt(14)
                        p.font.bold = True
                        p.level = 0
                    
                    # Add design elements
                    design = store.get('design_elements', '')
                    if design:
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = "Design: "
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run = p.add_run()
                        run.text = design
                        run.font.size = Pt(12)
                        p.level = 1
                    
                    # Add customer experience
                    experience = store.get('customer_experience', '')
                    if experience:
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = "Experience: "
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run = p.add_run()
                        run.text = experience
                        run.font.size = Pt(12)
                        p.level = 1
                    
                    # Add cultural meaning
                    cultural = store.get('cultural_meaning', '')
                    if cultural:
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = "Significance: "
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run = p.add_run()
                        run.text = cultural
                        run.font.size = Pt(12)
                        p.level = 1
                    
                    # Add spacing between stores
                    p = tf.add_paragraph()
                    p.text = ""
                else:
                    # Simple string store
                    p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                    first_paragraph = False
                    p.text = str(store)
                    p.font.size = Pt(13)
                    p.level = 0
        
        # Alternative experiences section (if no flagship stores)
        if alternative and not stores:
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            p.text = alternative
            p.font.size = Pt(14)
            p.level = 0

        # Safety net: if nothing was written, add a fallback so the slide isn't blank
        if first_paragraph:
            brand_name = data.get('title', '').replace('Flagship Stores & Experiences', '').strip()
            p = tf.paragraphs[0]
            p.text = "Flagship store and retail experience information is being compiled."
            p.font.size = Pt(14)
            p.level = 0
            logger.warning(f"Flagship stores slide had no content to render — fallback used")

        self._apply_body_font(tf)
        logger.info("Flagship stores slide created (using Title and Content layout)")
    
    def _create_collection_introduction_slide(self, data: Dict[str, Any]):
        """
        Create collection introduction slide using Title and Content layout.

        Args:
            data: Slide data containing title, content with collection_story, key_themes, highlights
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})

        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Collection Introduction')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Collection story section
        story = content.get('collection_story', {})
        if isinstance(story, dict):
            narrative = story.get('narrative', '')
            creative_direction = story.get('creative_direction', '')

            if narrative:
                p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
                first_paragraph = False
                p.text = narrative
                p.font.size = Pt(14)
                p.level = 0

            if creative_direction:
                p = tf.add_paragraph()
                p.text = ""  # Spacing

                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Creative Direction: "
                run.font.bold = True
                run.font.size = Pt(14)
                run = p.add_run()
                run.text = creative_direction
                run.font.size = Pt(14)
                p.level = 0

        # Key themes section
        themes = content.get('key_themes', {})
        if themes:
            color_palette = themes.get('color_palette', '')
            prints = themes.get('prints_and_patterns', '')
            silhouettes = themes.get('silhouettes', '')

            if color_palette or prints or silhouettes:
                p = tf.add_paragraph()
                p.text = ""  # Spacing

            if color_palette:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Color Palette: "
                run.font.bold = True
                run.font.size = Pt(14)
                run = p.add_run()
                run.text = color_palette
                run.font.size = Pt(14)
                p.level = 0

            if prints:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Prints & Patterns: "
                run.font.bold = True
                run.font.size = Pt(14)
                run = p.add_run()
                run.text = prints
                run.font.size = Pt(14)
                p.level = 0

            if silhouettes:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "Silhouettes: "
                run.font.bold = True
                run.font.size = Pt(14)
                run = p.add_run()
                run.text = silhouettes
                run.font.size = Pt(14)
                p.level = 0

        # Highlights section
        highlights = content.get('highlights', [])
        if highlights:
            p = tf.add_paragraph()
            p.text = ""  # Spacing

            p = tf.add_paragraph()
            p.text = "Collection Highlights:"
            p.font.bold = True
            p.font.size = Pt(14)
            p.level = 0

            for highlight in highlights:
                p = tf.add_paragraph()
                p.text = highlight
                p.font.size = Pt(13)
                p.level = 1

        # Fallback if nothing rendered
        if first_paragraph:
            p = tf.paragraphs[0]
            p.text = "Collection introduction details are being compiled."
            p.font.size = Pt(14)
            p.level = 0
            logger.warning("Collection introduction slide had no content — fallback used")

        self._apply_body_font(tf)
        logger.info("Collection introduction slide created (using Title and Content layout)")

    def _create_core_collection_slide(self, data: Dict[str, Any]):
        """
        Create core collection slide using Title and Content layout.

        Args:
            data: Slide data containing title, content with signature_categories (headline, description, iconic_staple)
        """
        slide = self.prs.slides.add_slide(self.title_content_layout)
        content = data.get('content', {})
        
        # Use title placeholder
        title = slide.shapes.title
        title.text = data.get('title', 'Core Collection & Signature Categories')
        self._apply_typo(title.text_frame.paragraphs[0].font, 'heading')

        # Use content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame

        first_paragraph = True

        # Add signature categories - LLM returns {headline, description, iconic_staple}
        categories = content.get('signature_categories', [])
        if categories:
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            p.text = "Signature Categories:"
            p.font.size = Pt(14)
            p.font.bold = True
            p.level = 0
            
            for category in categories:
                p = tf.add_paragraph()
                
                if isinstance(category, dict):
                    # Get headline (primary), fall back to category/name
                    cat_name = category.get('headline', category.get('category', category.get('name', '')))
                    cat_desc = category.get('description', '')
                    iconic_staple = category.get('iconic_staple')
                    
                    p.level = 1
                    
                    if cat_desc:
                        # Add category name in bold
                        run = p.add_run()
                        run.text = f"{cat_name}: "
                        run.font.bold = True
                        run.font.size = Pt(13)
                        
                        # Add description in normal text
                        run = p.add_run()
                        run.text = cat_desc
                        run.font.size = Pt(13)
                    else:
                        # Just the category name
                        p.text = cat_name
                        p.font.size = Pt(13)
                    
                    # Add iconic staple if available and not null
                    if iconic_staple and iconic_staple != 'null':
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = "Iconic Staple: "
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run = p.add_run()
                        run.text = str(iconic_staple)
                        run.font.size = Pt(12)
                        p.level = 2
                else:
                    # Simple string category
                    p.text = str(category)
                    p.font.size = Pt(13)
                    p.level = 1
        
        self._apply_body_font(tf)
        logger.info("Core collection slide created (using Title and Content layout)")
    
    async def _fetch_collection_items(self, collection_id: str) -> list:
        """
        Fetch all items for a collection from Firestore.
        
        Args:
            collection_id: ID of the collection
            
        Returns:
            List of item dictionaries
        """
        try:
            # Items are stored as a subcollection under each collection document
            items_ref = self.db.collection('collections').document(collection_id).collection('items')
            items_docs = items_ref.stream()
            
            items = []
            for doc in items_docs:
                item_data = doc.to_dict()
                item_data['item_id'] = doc.id
                items.append(item_data)
            
            # Filter out excluded items (included defaults to True if not set)
            items = [item for item in items if item.get('included', True) is not False]

            logger.info(f"Fetched {len(items)} included items for collection {collection_id}")
            return items
            
        except Exception as e:
            logger.error(f"Error fetching items: {e}")
            return []
    
    async def _generate_product_slides(self, collection_id: str, products_per_slide: int, category_order: dict = None):
        """
        Generate product slides based on collection items.

        Args:
            collection_id: ID of the collection
            products_per_slide: Number of products per slide (1, 2, or 4)
            category_order: Dict mapping category name -> display_order
        """
        if category_order is None:
            category_order = {}

        logger.info(f"Generating product slides ({products_per_slide} per slide)...")

        # Fetch items
        items = await self._fetch_collection_items(collection_id)

        if not items:
            logger.warning("No items found for product slides")
            return

        # Sort by category (alphabetical, nulls last), then subcategory
        # (nulls last), then display_order
        items.sort(key=lambda x: (
            x.get('category') or 'zzz',
            x.get('subcategory') or 'zzz',
            x.get('display_order', 0)
        ))

        # Group items by (category, subcategory) — only items sharing both
        # category AND subcategory will appear on the same slide
        items_by_cat_subcat = OrderedDict()
        for item in items:
            category = item.get('category', 'Uncategorized')
            subcategory = item.get('subcategory') or ''
            key = (category, subcategory)
            if key not in items_by_cat_subcat:
                items_by_cat_subcat[key] = []
            items_by_cat_subcat[key].append(item)

        unique_categories = list(dict.fromkeys(cat for cat, _ in items_by_cat_subcat))
        logger.info(
            f"Grouped {len(items)} items into {len(items_by_cat_subcat)} "
            f"subcategory groups across {len(unique_categories)} categories"
        )

        # Batch-translate categories, subcategories, and sales talk if non-English
        if self._language != "en":
            translatable = []
            for item in items:
                if item.get('category'):
                    translatable.append(item['category'])
                if item.get('subcategory'):
                    translatable.append(item['subcategory'])
                if item.get('sales_talk'):
                    translatable.append(item['sales_talk'].strip())
            self._translations = await self._translate_strings(translatable)

        # Sort groups: by category display_order first, then subcategory alpha
        sorted_keys = sorted(
            items_by_cat_subcat.keys(),
            key=lambda k: (category_order.get(k[0], 999), k[0], k[1] or 'zzz')
        )

        # Generate slides based on layout
        slides_generated = 0
        last_category = None

        for category, subcategory in sorted_keys:
            group_items = items_by_cat_subcat[(category, subcategory)]
            subcat_label = f" - {subcategory}" if subcategory else ""
            logger.info(
                f"Generating slides for {category}{subcat_label} "
                f"({len(group_items)} items)"
            )

            # Create category divider slide once per category
            if category != last_category:
                self._create_category_divider_slide(category)
                slides_generated += 1
                last_category = category

            # Create product slides — only matching items on each slide
            if products_per_slide == 1:
                for item in group_items:
                    self._create_1up_product_slide(item)
                    slides_generated += 1
            elif products_per_slide == 2:
                for i in range(0, len(group_items), 2):
                    items_group = group_items[i:i+2]
                    self._create_2up_product_slide(items_group)
                    slides_generated += 1
            elif products_per_slide == 3:
                for i in range(0, len(group_items), 3):
                    items_group = group_items[i:i+3]
                    self._create_3up_product_slide(items_group)
                    slides_generated += 1
            elif products_per_slide == 4:
                for i in range(0, len(group_items), 4):
                    items_group = group_items[i:i+4]
                    self._create_4up_product_slide(items_group)
                    slides_generated += 1

        logger.info(f"Generated {slides_generated} product slides")
    
    def _create_category_divider_slide(self, category: str):
        """
        Create a divider slide for a product category.
        
        Args:
            category: Category name
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        display_category = self._t(category)

        # Category title - centered
        title_box = slide.shapes.add_textbox(
            left=Inches(1 * self._w_scale),
            top=Inches(3.5),
            width=Inches(8 * self._w_scale),
            height=Inches(1)
        )
        tf = title_box.text_frame
        tf.text = display_category

        # Format title
        p = tf.paragraphs[0]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        self._apply_typo(p.font, 'heading', default_size=48)

        logger.info(f"Category divider slide created: {display_category}")
    
    def _download_image_as_stream(self, image_dict: dict) -> Optional[BytesIO]:
        """
        Download image from Firebase Storage or URL and return as BytesIO stream.
        
        Args:
            image_dict: Image dictionary with 'storage_path' and/or 'url'
            
        Returns:
            BytesIO stream of image data, or None if download fails
        """
        try:
            storage_path = image_dict.get('storage_path')
            image_url = image_dict.get('url')
            
            # Try storage path first (preferred)
            if storage_path:
                try:
                    logger.debug(f"Downloading image from storage path: {storage_path}")
                    blob = self.bucket.blob(storage_path)
                    image_bytes = blob.download_as_bytes()
                    return BytesIO(image_bytes)
                except Exception as e:
                    logger.warning(f"Failed to download from storage path: {e}")
                    # Fall through to URL fallback
            
            # Fallback to URL
            if image_url:
                try:
                    logger.debug(f"Downloading image from URL: {image_url}")
                    response = requests.get(image_url, timeout=10)
                    response.raise_for_status()
                    return BytesIO(response.content)
                except Exception as e:
                    logger.warning(f"Failed to download from URL: {e}")
            
            logger.warning("No valid storage_path or url found in image dict")
            return None
            
        except Exception as e:
            logger.error(f"Error downloading image: {e}")
            return None
    
    def _add_slide_title(self, slide, items):
        """
        Add a 'Category - SubCategory' title to the top-left of a product slide.

        Args:
            slide: The slide object to add the title to
            items: A single item dict or list of item dicts on this slide
        """
        if isinstance(items, dict):
            items = [items]
        if not items:
            return

        first = items[0]
        category = self._t(first.get('category', ''))
        subcategory = self._t(first.get('subcategory', ''))

        if category and subcategory:
            title_text = f"{category} - {subcategory}"
        elif category:
            title_text = category
        else:
            return

        title_box = slide.shapes.add_textbox(
            left=Inches(0.4),
            top=Inches(0.15),
            width=Inches(8),
            height=Inches(0.6)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        self._apply_typo(p.font, 'heading', default_size=18)

    def _add_highlight_star(self, slide, image_left, image_top):
        """
        Add a grey star icon to the top-left of a product image for highlighted items.

        Args:
            slide: The slide object
            image_left: Left position of the product image (Emu)
            image_top: Top position of the product image (Emu)
        """
        star_size = Inches(0.3)
        star_box = slide.shapes.add_textbox(
            left=image_left - Inches(0.35),
            top=image_top,
            width=star_size,
            height=star_size
        )
        tf = star_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "\u2605"  # Black star
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)  # Neutral grey
        p.alignment = PP_ALIGN.CENTER

    def _add_sales_talk_below(self, slide, tf, item: dict, details_left,
                              details_top, image_width, font_size: int = 12):
        """
        Add sales talk in a separate text box positioned right below the
        existing details content. Width matches the image.

        Estimates the rendered height of paragraphs in the details text
        frame to calculate the vertical position.

        Args:
            slide: The slide object
            tf: The details text frame (used to count paragraphs)
            item: Item dictionary
            details_left: Left position of the details text box
            details_top: Top position of the details text box (in inches)
            image_width: Width to constrain the sales talk to
            font_size: Font size in points
        """
        sales_talk = (item.get('sales_talk') or '').strip()
        if not sales_talk:
            return
        sales_talk = self._t(sales_talk)

        # Estimate content height from paragraphs in the text frame
        total_height_pt = 0
        for p in tf.paragraphs:
            if p.font.size:
                total_height_pt += p.font.size.pt * 1.3
            else:
                total_height_pt += 8  # empty/spacing paragraph
        content_height_inches = total_height_pt / 72.0

        sales_top = details_top + content_height_inches

        box = slide.shapes.add_textbox(
            left=details_left,
            top=Inches(sales_top),
            width=image_width,
            height=Inches(0.5)
        )
        stf = box.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        p.text = sales_talk
        p.font.bold = True
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.CENTER
        self._apply_body_font(stf)

    def _create_1up_product_slide(self, item: dict):
        """
        Create a slide with 1 product (full slide layout).

        Layout:
        - Product image: Left side (if available)
        - Product details: Right side

        Args:
            item: Item dictionary
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_slide_title(slide, item)

        # Extract item data
        product_name = item.get('product_name', 'Unknown Product')
        sku = item.get('sku', '')
        color = item.get('color', '')
        subcategory = item.get('subcategory', '')
        origin = item.get('origin', '')
        description = item.get('description', '')
        materials = item.get('materials', [])
        wholesale_price = item.get('wholesale_price')
        rrp = item.get('rrp')
        currency = item.get('currency', 'USD')
        images = item.get('images', [])

        # Left side: Product image or placeholder
        image_added = False
        
        if images and len(images) > 0:
            # Try to download and embed the first image
            image_stream = self._download_image_as_stream(images[0])
            
            if image_stream:
                try:
                    # Add image to slide with size constraints
                    # Smaller width (3.5" instead of 4.5") and max height (5")
                    picture = slide.shapes.add_picture(
                        image_stream,
                        left=Inches(0.75 * self._w_scale),
                        top=Inches(1.5),
                        width=Inches(3.5),
                        height=Inches(5)
                    )
                    image_added = True
                    logger.debug(f"Image added for product: {product_name}")
                except Exception as e:
                    logger.warning(f"Failed to add image to slide: {e}")

        # Add highlight star if item is highlighted
        if item.get('highlighted_item'):
            self._add_highlight_star(slide, Inches(0.75 * self._w_scale), Inches(1.5))

        # Show placeholder if no image was added
        if not image_added:
            image_box = slide.shapes.add_textbox(
                left=Inches(0.5 * self._w_scale),
                top=Inches(1),
                width=Inches(4.5),
                height=Inches(6)
            )
            tf = image_box.text_frame
            
            if images and len(images) > 0:
                tf.text = "[Image Not Available]"
            else:
                tf.text = "[No Image]"
            
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
        
        # Right side: Product details
        details_box = slide.shapes.add_textbox(
            left=Inches(5.5 * self._w_scale),
            top=Inches(1),
            width=Inches(4 * self._w_scale),
            height=Inches(6)
        )
        tf = details_box.text_frame
        
        # Product name
        if self._show_fields.get('product_name', True):
            p = tf.paragraphs[0]
            p.text = product_name
            p.font.bold = True
            self._apply_typo(p.font, 'heading', default_size=24)

        # SKU
        if sku and self._show_fields.get('sku', True):
            p = tf.add_paragraph()
            p.text = f"SKU: {sku}"
            p.font.size = Pt(12)
            p.font.italic = True

        # Subcategory
        if subcategory:
            p = tf.add_paragraph()
            p.text = f"Subcategory: {subcategory}"
            p.font.size = Pt(13)

        # Color
        if color and self._show_fields.get('color', True):
            p = tf.add_paragraph()
            p.text = f"Color: {color}"
            p.font.size = Pt(14)

        # Origin
        if origin and self._show_fields.get('origin', True):
            p = tf.add_paragraph()
            p.text = f"Origin: {origin}"
            p.font.size = Pt(13)

        # Spacing
        p = tf.add_paragraph()
        p.text = ""

        # Description
        if description and self._show_fields.get('description', True):
            p = tf.add_paragraph()
            p.text = description
            p.font.size = Pt(12)

        # Spacing
        p = tf.add_paragraph()
        p.text = ""

        # Materials
        if materials and self._show_fields.get('material', True):
            p = tf.add_paragraph()
            p.text = "Materials:"
            p.font.size = Pt(12)
            p.font.bold = True

            for material in materials:
                p = tf.add_paragraph()
                p.text = f"• {material}"
                p.font.size = Pt(11)

        # Pricing
        show_wp = self._show_fields.get('wholesale_price', False)
        show_rrp = self._show_fields.get('rrp', True)
        wp_val = wholesale_price if show_wp else None
        rrp_val = rrp if show_rrp else None
        if wp_val or rrp_val:
            p = tf.add_paragraph()
            p.text = ""

            p = tf.add_paragraph()
            pricing_text = ""
            if wp_val:
                pricing_text += f"Wholesale: {currency} {wholesale_price:.2f}"
            if rrp_val:
                if pricing_text:
                    pricing_text += " | "
                pricing_text += f"RRP: {currency} {rrp:.2f}"
            p.text = pricing_text
            p.font.size = Pt(11)
        
        self._apply_body_font(tf)

        # Sales talk — fixed near bottom of the details column (1-up only)
        sales_talk = (item.get('sales_talk') or '').strip()
        if sales_talk:
            sales_talk = self._t(sales_talk)
            box = slide.shapes.add_textbox(
                left=Inches(5.5 * self._w_scale),
                top=Inches(6.0),
                width=Inches(4 * self._w_scale),
                height=Inches(0.8)
            )
            stf = box.text_frame
            stf.word_wrap = True
            p = stf.paragraphs[0]
            p.text = sales_talk
            p.font.bold = True
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            self._apply_body_font(stf)

        logger.info(f"1-up product slide created: {product_name}")

    def _create_2up_product_slide(self, items: list):
        """
        Create a slide with 2 products (side by side).
        
        Layout: Two columns, each with image on left and details on right.
        
        Args:
            items: List of 1-2 item dictionaries
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_slide_title(slide, items)

        # Process up to 2 items
        for idx, item in enumerate(items[:2]):
            # Calculate position based on column (0 = left, 1 = right)
            column_offset = idx * 5 * self._w_scale  # 5 inches between columns, scaled

            # Extract item data
            product_name = item.get('product_name', 'Unknown Product')
            sku = item.get('sku', '')
            color = item.get('color', '')
            subcategory = item.get('subcategory', '')
            origin = item.get('origin', '')
            description = item.get('description', '')
            materials = item.get('materials', [])
            wholesale_price = item.get('wholesale_price')
            rrp = item.get('rrp')
            currency = item.get('currency', 'USD')
            images = item.get('images', [])
            
            # Product image or placeholder (left side of column)
            image_added = False
            image_left = Inches(0.5 * self._w_scale + column_offset)
            image_top = Inches(1.5)
            image_width = Inches(2)
            image_height = Inches(3)
            
            if images and len(images) > 0:
                image_stream = self._download_image_as_stream(images[0])
                
                if image_stream:
                    try:
                        slide.shapes.add_picture(
                            image_stream,
                            left=image_left,
                            top=image_top,
                            width=image_width,
                            height=image_height
                        )
                        image_added = True
                    except Exception as e:
                        logger.warning(f"Failed to add image: {e}")

            # Add highlight star if item is highlighted
            if item.get('highlighted_item'):
                self._add_highlight_star(slide, image_left, image_top)

            # Show placeholder if no image
            if not image_added:
                image_box = slide.shapes.add_textbox(
                    left=image_left,
                    top=image_top,
                    width=image_width,
                    height=image_height
                )
                tf = image_box.text_frame
                tf.text = "[No Image]" if not images else "[Image Not Available]"
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER

            # Product details (right side of column)
            details_left = Inches(2.75 * self._w_scale + column_offset)
            details_width = Inches(2 * self._w_scale)
            
            details_box = slide.shapes.add_textbox(
                left=details_left,
                top=Inches(1),
                width=details_width,
                height=Inches(6)
            )
            tf = details_box.text_frame
            tf.word_wrap = True
            
            # Product name
            if self._show_fields.get('product_name', True):
                p = tf.paragraphs[0]
                p.text = product_name
                p.font.bold = True
                self._apply_typo(p.font, 'heading', default_size=16)

            # SKU
            if sku and self._show_fields.get('sku', True):
                p = tf.add_paragraph()
                p.text = f"SKU: {sku}"
                p.font.size = Pt(9)
                p.font.italic = True

            # Subcategory
            if subcategory:
                p = tf.add_paragraph()
                p.text = f"Subcat: {subcategory}"
                p.font.size = Pt(9)

            # Color
            if color and self._show_fields.get('color', True):
                p = tf.add_paragraph()
                p.text = f"Color: {color}"
                p.font.size = Pt(10)

            # Origin
            if origin and self._show_fields.get('origin', True):
                p = tf.add_paragraph()
                p.text = f"Origin: {origin}"
                p.font.size = Pt(9)

            # Spacing
            p = tf.add_paragraph()
            p.text = ""

            # Description (truncate if too long)
            if description and self._show_fields.get('description', True):
                p = tf.add_paragraph()
                # Limit description to 150 chars for 2-up layout
                if len(description) > 150:
                    p.text = description[:147] + "..."
                else:
                    p.text = description
                p.font.size = Pt(9)

            # Materials (condensed)
            if materials and self._show_fields.get('material', True):
                p = tf.add_paragraph()
                p.text = "Materials:"
                p.font.size = Pt(9)
                p.font.bold = True

                # Show first 3 materials
                for material in materials[:3]:
                    p = tf.add_paragraph()
                    p.text = f"• {material}"
                    p.font.size = Pt(8)

                if len(materials) > 3:
                    p = tf.add_paragraph()
                    p.text = f"+ {len(materials) - 3} more"
                    p.font.size = Pt(8)
                    p.font.italic = True

            # Pricing
            show_wp = self._show_fields.get('wholesale_price', False)
            show_rrp = self._show_fields.get('rrp', True)
            wp_val = wholesale_price if show_wp else None
            rrp_val = rrp if show_rrp else None
            if wp_val or rrp_val:
                p = tf.add_paragraph()
                p.text = ""

                p = tf.add_paragraph()
                pricing_text = ""
                if wp_val:
                    pricing_text += f"W: {currency} {wholesale_price:.2f}"
                if rrp_val:
                    if pricing_text:
                        pricing_text += " | "
                    pricing_text += f"RRP: {currency} {rrp:.2f}"
                p.text = pricing_text
                p.font.size = Pt(9)
        
            self._apply_body_font(tf)

            self._add_sales_talk_below(
                slide, tf, item,
                details_left=details_left,
                details_top=1.0,
                image_width=image_width,
                font_size=12
            )

        logger.info(f"2-up product slide created with {len(items)} item(s)")

    def _create_3up_product_slide(self, items: list):
        """
        Create a slide with 3 products (line sheet style).
        
        Layout: Images in top row, details in bottom row (3 columns).
        
        Args:
            items: List of 1-3 item dictionaries
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_slide_title(slide, items)

        # Column positions (left edge for each product) - 3 columns with more space
        column_positions = [1.0 * self._w_scale, 4.0 * self._w_scale, 7.0 * self._w_scale]
        column_width = 2.5 * self._w_scale
        
        # Process up to 3 items
        for idx, item in enumerate(items[:3]):
            if idx >= len(column_positions):
                break
                
            col_left = Inches(column_positions[idx])
            
            # Extract item data
            product_name = item.get('product_name', 'Unknown Product')
            sku = item.get('sku', '')
            color = item.get('color', '')
            subcategory = item.get('subcategory', '')
            origin = item.get('origin', '')
            description = item.get('description', '')
            materials = item.get('materials', [])
            wholesale_price = item.get('wholesale_price')
            rrp = item.get('rrp')
            currency = item.get('currency', 'USD')
            images = item.get('images', [])
            
            # Top row: Product image (keep original size, don't stretch)
            image_added = False
            image_width = Inches(2.5)  # Original 3-up image width
            image_height = Inches(3.5)
            image_top = Inches(0.75)

            if images and len(images) > 0:
                image_stream = self._download_image_as_stream(images[0])

                if image_stream:
                    try:
                        slide.shapes.add_picture(
                            image_stream,
                            left=col_left,
                            top=image_top,
                            width=image_width,
                            height=image_height
                        )
                        image_added = True
                    except Exception as e:
                        logger.warning(f"Failed to add image: {e}")

            # Add highlight star if item is highlighted
            if item.get('highlighted_item'):
                self._add_highlight_star(slide, col_left, image_top)

            # Show placeholder if no image
            if not image_added:
                image_box = slide.shapes.add_textbox(
                    left=col_left,
                    top=image_top,
                    width=image_width,
                    height=image_height
                )
                tf = image_box.text_frame
                tf.text = "[No Image]" if not images else "[N/A]"
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER

            # Bottom row: Product details
            details_top = Inches(4.5)
            details_height = Inches(3)

            details_box = slide.shapes.add_textbox(
                left=col_left,
                top=details_top,
                width=Inches(column_width),
                height=details_height
            )
            tf = details_box.text_frame
            tf.word_wrap = True
            
            # Product name
            if self._show_fields.get('product_name', True):
                p = tf.paragraphs[0]
                p.text = product_name
                p.font.bold = True
                self._apply_typo(p.font, 'heading', default_size=11)

            # SKU
            if sku and self._show_fields.get('sku', True):
                p = tf.add_paragraph()
                p.text = f"SKU: {sku}"
                p.font.size = Pt(8)
                p.font.italic = True

            # Color
            if color and self._show_fields.get('color', True):
                p = tf.add_paragraph()
                p.text = f"Color: {color}"
                p.font.size = Pt(8)

            # Subcategory
            if subcategory:
                p = tf.add_paragraph()
                p.text = f"Subcat: {subcategory}"
                p.font.size = Pt(8)

            # Origin
            if origin and self._show_fields.get('origin', True):
                p = tf.add_paragraph()
                p.text = f"Origin: {origin}"
                p.font.size = Pt(8)

            # Description (dynamic font sizing, word wrap enabled)
            if description and self._show_fields.get('description', True):
                p = tf.add_paragraph()
                p.text = description

                # Dynamic font size based on length
                if len(description) <= 80:
                    p.font.size = Pt(8)
                elif len(description) <= 120:
                    p.font.size = Pt(7)
                else:
                    p.font.size = Pt(6)

            # Materials (show first 4)
            if materials and self._show_fields.get('material', True):
                p = tf.add_paragraph()
                p.text = "Materials:"
                p.font.size = Pt(8)
                p.font.bold = True

                # Show first 4 materials
                for material in materials[:4]:
                    p = tf.add_paragraph()
                    p.text = f"• {material}"
                    p.font.size = Pt(7)

                # If more than 4, show count
                if len(materials) > 4:
                    p = tf.add_paragraph()
                    p.text = f"+ {len(materials) - 4} more"
                    p.font.size = Pt(7)
                    p.font.italic = True

            # Pricing
            if wholesale_price and self._show_fields.get('wholesale_price', False):
                p = tf.add_paragraph()
                p.text = f"W: {currency} {wholesale_price:.2f}"
                p.font.size = Pt(8)

            if rrp and self._show_fields.get('rrp', True):
                p = tf.add_paragraph()
                p.text = f"RRP: {currency} {rrp:.2f}"
                p.font.size = Pt(8)
        
            self._apply_body_font(tf)

            self._add_sales_talk_below(
                slide, tf, item,
                details_left=col_left,
                details_top=4.5,
                image_width=Inches(2.5),
                font_size=10
            )

        logger.info(f"3-up product slide created with {len(items)} item(s)")

    def _create_4up_product_slide(self, items: list):
        """
        Create a slide with 4 products (line sheet style).
        
        Layout: Images in top row, details in bottom row (4 columns).
        
        Args:
            items: List of 1-4 item dictionaries
        """
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_slide_title(slide, items)

        # Column positions (left edge for each product)
        column_positions = [0.5 * self._w_scale, 2.9 * self._w_scale, 5.3 * self._w_scale, 7.7 * self._w_scale]
        column_width = 2.1 * self._w_scale
        
        # Process up to 4 items
        for idx, item in enumerate(items[:4]):
            if idx >= len(column_positions):
                break
                
            col_left = Inches(column_positions[idx])
            
            # Extract item data
            product_name = item.get('product_name', 'Unknown Product')
            sku = item.get('sku', '')
            color = item.get('color', '')
            subcategory = item.get('subcategory', '')
            origin = item.get('origin', '')
            description = item.get('description', '')
            materials = item.get('materials', [])
            wholesale_price = item.get('wholesale_price')
            rrp = item.get('rrp')
            currency = item.get('currency', 'USD')
            images = item.get('images', [])
            
            # Top row: Product image
            image_added = False
            image_width = Inches(1.8)
            image_height = Inches(2.5)
            image_top = Inches(1)
            
            if images and len(images) > 0:
                image_stream = self._download_image_as_stream(images[0])
                
                if image_stream:
                    try:
                        slide.shapes.add_picture(
                            image_stream,
                            left=col_left,
                            top=image_top,
                            width=image_width,
                            height=image_height
                        )
                        image_added = True
                    except Exception as e:
                        logger.warning(f"Failed to add image: {e}")

            # Add highlight star if item is highlighted
            if item.get('highlighted_item'):
                self._add_highlight_star(slide, col_left, image_top)

            # Show placeholder if no image
            if not image_added:
                image_box = slide.shapes.add_textbox(
                    left=col_left,
                    top=image_top,
                    width=image_width,
                    height=image_height
                )
                tf = image_box.text_frame
                tf.text = "[No Image]" if not images else "[N/A]"
                p = tf.paragraphs[0]
                p.font.size = Pt(10)
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER

            # Bottom row: Product details
            details_top = Inches(4)
            details_height = Inches(3.5)
            
            details_box = slide.shapes.add_textbox(
                left=col_left,
                top=details_top,
                width=Inches(column_width),
                height=details_height
            )
            tf = details_box.text_frame
            tf.word_wrap = True
            
            # Product name
            if self._show_fields.get('product_name', True):
                p = tf.paragraphs[0]
                p.text = product_name
                p.font.bold = True
                self._apply_typo(p.font, 'heading', default_size=9)

            # SKU
            if sku and self._show_fields.get('sku', True):
                p = tf.add_paragraph()
                p.text = f"SKU: {sku}"
                p.font.size = Pt(7)
                p.font.italic = True

            # Color
            if color and self._show_fields.get('color', True):
                p = tf.add_paragraph()
                p.text = f"Color: {color}"
                p.font.size = Pt(7)

            # Subcategory
            if subcategory:
                p = tf.add_paragraph()
                p.text = f"Subcat: {subcategory}"
                p.font.size = Pt(7)

            # Origin
            if origin and self._show_fields.get('origin', True):
                p = tf.add_paragraph()
                p.text = f"Origin: {origin}"
                p.font.size = Pt(7)

            # Description (dynamic font sizing, word wrap enabled)
            if description and self._show_fields.get('description', True):
                p = tf.add_paragraph()
                p.text = description

                # Dynamic font size based on length
                if len(description) <= 60:
                    p.font.size = Pt(7)
                elif len(description) <= 100:
                    p.font.size = Pt(6)
                else:
                    p.font.size = Pt(5)

            # Materials (show first 4)
            if materials and self._show_fields.get('material', True):
                p = tf.add_paragraph()
                p.text = "Materials:"
                p.font.size = Pt(7)
                p.font.bold = True

                # Show first 4 materials
                for material in materials[:4]:
                    p = tf.add_paragraph()
                    p.text = f"• {material}"
                    p.font.size = Pt(6)

                # If more than 4, show count
                if len(materials) > 4:
                    p = tf.add_paragraph()
                    p.text = f"+ {len(materials) - 4} more"
                    p.font.size = Pt(6)
                    p.font.italic = True

            # Pricing
            if wholesale_price and self._show_fields.get('wholesale_price', False):
                p = tf.add_paragraph()
                p.text = f"W: {currency} {wholesale_price:.2f}"
                p.font.size = Pt(7)

            if rrp and self._show_fields.get('rrp', True):
                p = tf.add_paragraph()
                p.text = f"RRP: {currency} {rrp:.2f}"
                p.font.size = Pt(7)
        
            self._apply_body_font(tf)

            self._add_sales_talk_below(
                slide, tf, item,
                details_left=col_left,
                details_top=4.0,
                image_width=Inches(1.8),
                font_size=9
            )

        logger.info(f"4-up product slide created with {len(items)} item(s)")

    async def _save_and_upload(self, collection_id: str) -> str:
        """
        Save presentation to temp file, upload to Firebase Storage, and clean up.
        
        Args:
            collection_id: ID of the collection
            
        Returns:
            Public download URL for the presentation
        """
        try:
            # Generate temp file path
            temp_dir = tempfile.gettempdir()
            timestamp = datetime.now().timestamp()
            temp_path = os.path.join(temp_dir, f"{collection_id}_{timestamp}.pptx")
            
            logger.info(f"Saving presentation to: {temp_path}")
            
            # Save presentation
            self.prs.save(temp_path)

            # Upload to Firebase Storage with unique path per generation
            # Using a unique filename prevents GCS/CDN from serving cached old versions
            gen_ts = int(datetime.now().timestamp())
            storage_path = f"presentations/{collection_id}/{gen_ts}_presentation.pptx"
            blob = self.bucket.blob(storage_path)
            blob.cache_control = 'no-cache, no-store, must-revalidate'

            logger.info(f"Uploading to Firebase Storage: {blob.name}")
            blob.upload_from_filename(
                temp_path,
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )

            # Make public and get URL
            blob.make_public()
            download_url = blob.public_url

            # Clean up temp file
            os.remove(temp_path)
            logger.info(f"Temp file cleaned up: {temp_path}")
            collection_ref = self.db.collection('collections').document(collection_id)
            collection_ref.update({
                'presentation': {
                    'generated_at': datetime.utcnow(),
                    'download_url': download_url,
                    'slide_count': len(self.prs.slides),
                    'storage_path': storage_path
                }
            })
            
            logger.info(f"Presentation metadata saved to Firestore")
            
            return download_url
            
        except Exception as e:
            logger.error(f"Error saving and uploading presentation: {e}")
            # Clean up temp file if it exists
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise


# Global service instance
presentation_generation_service = PresentationGenerationService()
